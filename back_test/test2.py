from transformers import pipeline
from sentence_transformers import SentenceTransformer
from faiss import IndexFlatIP
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
from pptx import Presentation
import chardet
import os

# Шаг 1: Индексация документов

def detect_encoding(file_path):
    with open(file_path, 'rb') as file:
        raw_data = file.read()
    result = chardet.detect(raw_data)
    return result['encoding']

def extract_text_from_pdf(file_path):
    doc = fitz.open(file_path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text

def extract_text_from_image(file_path):
    image = Image.open(file_path)
    text = pytesseract.image_to_string(image)
    return text

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

def extract_text_from_document(file_path):
    if file_path.endswith('.pdf'):
        return extract_text_from_pdf(file_path)
    elif file_path.endswith('.jpg') or file_path.endswith('.png'):
        return extract_text_from_image(file_path)
    elif file_path.endswith('.pptx'):
        return extract_text_from_pptx(file_path)
    else:
        try:
            encoding = detect_encoding(file_path)
            with open(file_path, 'r', encoding=encoding) as file:
                return file.read()
        except UnicodeDecodeError:
            # Если все еще возникает ошибка, попробуем использовать кодировку по умолчанию
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                return file.read()

# Загрузка модели эмбеддингов
embedding_model = SentenceTransformer('all-MiniLM-L6-v2')

# Получение списка файлов из директории
rawDocuments = os.listdir("./documents")
rawDocuments = [f"./documents/{filePath}" for filePath in rawDocuments]

# Представление документов в виде векторов
documents = [extract_text_from_document(file_path) for file_path in rawDocuments]
document_embeddings = embedding_model.encode(documents)

# Создание индекса FAISS
index = IndexFlatIP(document_embeddings.shape[1])
index.add(document_embeddings)

# Шаг 2: Поиск релевантных документов

# Представление запроса в виде вектора
query = "Запрос пользователя"
query_embedding = embedding_model.encode([query])

# Поиск ближайших соседей
distances, indices = index.search(query_embedding, k=2)
relevant_documents = [documents[i] for i in indices[0]]

# Шаг 3: Генерация ответа

# Подготовка контекста
context = "\n".join(relevant_documents)

# Загрузка модели генерации
generator = pipeline('text-generation', model='ColPali')

# Генерация ответа
response = generator(f"Запрос: {query}\nКонтекст: {context}\nОтвет:")

print(response[0]['generated_text'])