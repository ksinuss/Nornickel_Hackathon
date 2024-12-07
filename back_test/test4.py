import os
import shutil
import chardet
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
from pptx import Presentation
from transformers import AutoTokenizer, AutoModelForCausalLM, Trainer, TrainingArguments, RagRetriever, RagTokenForGeneration
from datasets import Dataset
import torch

# Шаг 1: Загрузка и предобработка данных

def detect_encoding(file_path):
    with open(file_path, 'rb') as file:
        raw_data = file.read()
    result = chardet.detect(raw_data)
    return result['encoding']

def extract_text_from_pdf(file_path):
    doc = fitz.open(file_path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text

def extract_text_from_image(file_path):
    image = Image.open(file_path)
    text = pytesseract.image_to_string(image)
    return text

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

def extract_text_from_document(file_path):
    if file_path.endswith('.pdf'):
        return extract_text_from_pdf(file_path)
    elif file_path.endswith('.jpg') or file_path.endswith('.png'):
        return extract_text_from_image(file_path)
    elif file_path.endswith('.pptx'):
        return extract_text_from_pptx(file_path)
    else:
        try:
            encoding = detect_encoding(file_path)
            with open(file_path, 'r', encoding=encoding) as file:
                return file.read()
        except UnicodeDecodeError:
            # Если все еще возникает ошибка, попробуем использовать кодировку по умолчанию
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                return file.read()

# Получение списка файлов из директории
rawDocuments = os.listdir("./documents")
rawDocuments = [f"./documents/{filePath}" for filePath in rawDocuments]

# Извлечение текста из документов
documents = [extract_text_from_document(file_path) for file_path in rawDocuments]

# Шаг 2: Обучение модели

# Загрузка токенизатора и модели
tokenizer = AutoTokenizer.from_pretrained('distilgpt2')
model = AutoModelForCausalLM.from_pretrained('distilgpt2')

# Добавление pad_token в токенизатор и изменение размера эмбеддингов модели
if tokenizer.pad_token is None:
    tokenizer.add_special_tokens({'pad_token': tokenizer.eos_token})
    model.resize_token_embeddings(len(tokenizer))

# Предобработка данных
def preprocess_function(examples):
    return tokenizer(examples['text'], truncation=True, padding='max_length', max_length=512)

# Создание датасета
dataset = Dataset.from_dict({'text': documents})
tokenized_dataset = dataset.map(preprocess_function, batched=True)

# Определение параметров обучения
training_args = TrainingArguments(
    output_dir='./results',
    num_train_epochs=3,
    per_device_train_batch_size=2,
    save_steps=10_000,
    save_total_limit=2,
)

# Удаление директории для сохранения результатов, если она существует
if os.path.exists(training_args.output_dir):
    shutil.rmtree(training_args.output_dir)

# Создание директории для сохранения результатов
os.makedirs(training_args.output_dir, exist_ok=True)

# Создание собственного класса тренера с переопределенным методом compute_loss
class CustomTrainer(Trainer):
    def compute_loss(self, model, inputs, return_outputs=False, *args, **kwargs):
        outputs = model(**inputs)
        logits = outputs.logits
        labels = inputs['input_ids']
        loss_fct = torch.nn.CrossEntropyLoss()
        loss = loss_fct(logits.view(-1, logits.size(-1)), labels.view(-1))
        return (loss, outputs) if return_outputs else loss

# Создание тренера
trainer = CustomTrainer(
    model=model,
    args=training_args,
    train_dataset=tokenized_dataset,
)

# Обучение модели
trainer.train()

# Шаг 3: Сохранение обученной модели
model.save_pretrained('./trained_model')
tokenizer.save_pretrained('./trained_model')

# Шаг 4: Генерация ответов на примеры вопросов

# Примеры вопросов
questions = [
    "Долевое участие металлов в EBITDA компании? То есть какой % EBITDA и EBITDA margin принесли продажи Ni, Cu и т д",
    "Как эта доля менялась от года к году (наверно нужно приложить отчеты хотя бы за 2-3 года)?",
    "Топ 5 факторов снижения EBITDA г/г? Также и EBITDA margin",
    "Какой актив формирует наибольшую долю FCF? Какой актив на 2 месте?",
    "Какие активы являются убыточными? Можно ли их исключить из производственной цепочки без последствий для общего цикла производства металлов?",
    "Какой прогноз по содержанию металлов руд НН на следующие 5 лет?",
    "Какую долю рынка занимает НН в производстве металлов? Как изменялась доля НН г/г? Какой прогноз на 5 лет?",
    "Как изменялся прогноз по предложению и спросу на Ni, Cu, Au в отчетах за разные года, например какой был прогноз в отчете за 2020 год на 2021/22/23? Какое было фактическое предложение и спрос в годовых отчетах 2021/22/23? Исходя из этого сделай вывод о целесообразности доверия таким прогнозам в целом"
]

# Загрузка обученной модели
trained_model = AutoModelForCausalLM.from_pretrained('./trained_model')
trained_tokenizer = AutoTokenizer.from_pretrained('./trained_model')

# Генерация ответов
def generate_answer(question, model, tokenizer, max_new_tokens=100):
    inputs = tokenizer(question, return_tensors='pt', truncation=True, padding='max_length', max_length=512)
    outputs = model.generate(inputs['input_ids'], attention_mask=inputs['attention_mask'], max_new_tokens=max_new_tokens)
    answer = tokenizer.decode(outputs[0], skip_special_tokens=True)
    return answer

# Ответы на вопросы
for question in questions:
    answer = generate_answer(question, trained_model, trained_tokenizer, max_new_tokens=100)
    print(f"Вопрос: {question}\nОтвет: {answer}\n")

# Использование RAG
def generate_answer_with_rag(question, model, tokenizer, max_new_tokens=100):
    inputs = tokenizer(question, return_tensors='pt', truncation=True, padding='max_length', max_length=512)
    outputs = model.generate(inputs['input_ids'], attention_mask=inputs['attention_mask'], max_new_tokens=max_new_tokens)
    answer = tokenizer.decode(outputs[0], skip_special_tokens=True)
    return answer

# Загрузка RAG модели и ретривера с аргументом trust_remote_code=True
retriever = RagRetriever.from_pretrained('facebook/rag-token-nq', index_name="exact", use_dummy_dataset=True, trust_remote_code=True)
rag_model = RagTokenForGeneration.from_pretrained('facebook/rag-token-nq', retriever=retriever, trust_remote_code=True)

# Ответы на вопросы с использованием RAG
for question in questions:
    answer = generate_answer_with_rag(question, rag_model, trained_tokenizer, max_new_tokens=100)
    print(f"Вопрос (RAG): {question}\nОтвет: {answer}\n")